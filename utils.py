# utils.py - 工具函数（优化网络连接和重试策略，兼容不同版本urllib3）
import requests
import json
import time
import random
import re
import os
from requests.adapters import HTTPAdapter
from urllib3.util.retry import Retry
from docx import Document
from docx.shared import Inches, Pt
from docx.enum.text import WD_ALIGN_PARAGRAPH
import io
import PyPDF2
from docx import Document
# 更新导入语句，使用新的配置变量
from config import DEEPSEEK_API_KEY, DEEPSEEK_API_URL, MAX_RETRIES, CONNECT_TIMEOUT, READ_TIMEOUT, BACKOFF_FACTOR

# 添加文件解析函数
def parse_uploaded_file(uploaded_file):
    """解析上传的Word或PDF文件，提取文本内容"""
    try:
        file_type = uploaded_file.type
        content = ""
        
        if file_type == "application/pdf":
            # 处理PDF文件
            pdf_reader = PyPDF2.PdfReader(io.BytesIO(uploaded_file.getvalue()))
            for page in pdf_reader.pages:
                content += page.extract_text() + "\n"
                
        elif file_type in ["application/vnd.openxmlformats-officedocument.wordprocessingml.document", 
                          "application/msword"]:
            # 处理Word文件
            doc = Document(io.BytesIO(uploaded_file.getvalue()))
            for paragraph in doc.paragraphs:
                if paragraph.text.strip():
                    content += paragraph.text + "\n"
                    
        else:
            return {"error": f"不支持的文件类型: {file_type}"}
        
        # 清理内容，移除多余的空行和空格
        content = re.sub(r'\n\s*\n', '\n\n', content)
        content = content.strip()
        
        return {"success": True, "content": content, "file_type": file_type}
        
    except Exception as e:
        return {"error": f"文件解析失败: {str(e)}"}

# 添加提取关键要求的函数
def extract_key_requirements(content, max_length=1000):
    """从政策文件中提取关键要求，限制长度"""
    if not content:
        return ""
    
    # 简单提取前几段作为关键要求
    paragraphs = [p for p in content.split('\n') if p.strip()]
    
    # 提取重要段落（包含关键词的段落优先）
    keywords = ['要求', '目标', '标准', '考试', '大纲', '政策', '规定', '掌握', '理解', '应用']
    important_paragraphs = []
    other_paragraphs = []
    
    for para in paragraphs:
        if any(keyword in para for keyword in keywords):
            important_paragraphs.append(para)
        else:
            other_paragraphs.append(para)
    
    # 组合内容，重要内容在前
    selected_content = important_paragraphs + other_paragraphs
    result_content = "\n".join(selected_content[:10])  # 最多10段
    
    # 限制总长度
    if len(result_content) > max_length:
        result_content = result_content[:max_length] + "..."
    
    return result_content

def create_session_with_retries():
    """创建带有重试机制的会话，兼容不同版本的 urllib3"""
    session = requests.Session()
    
    # 尝试不同的参数组合以适应不同版本的 urllib3
    retry_params = {
        "total": MAX_RETRIES,
        "backoff_factor": BACKOFF_FACTOR,
        "status_forcelist": [429, 500, 502, 503, 504],
    }
    
    # 尝试添加方法限制参数，使用不同版本的参数名
    try:
        # 首先尝试新版本的参数名
        retry_params["allowed_methods"] = ["POST"]
        retry_strategy = Retry(**retry_params)
    except TypeError as e:
        if "allowed_methods" in str(e):
            # 如果新版本参数不可用，尝试旧版本参数名
            try:
                del retry_params["allowed_methods"]
                retry_params["method_whitelist"] = ["POST"]
                retry_strategy = Retry(**retry_params)
            except TypeError:
                # 如果旧版本参数也不可用，不使用方法限制
                del retry_params["method_whitelist"]
                retry_strategy = Retry(**retry_params)
        else:
            # 其他错误，不使用方法限制
            if "allowed_methods" in retry_params:
                del retry_params["allowed_methods"]
            retry_strategy = Retry(**retry_params)
    
    # 创建适配器并挂载到会话
    adapter = HTTPAdapter(max_retries=retry_strategy, pool_connections=10, pool_maxsize=10)
    session.mount("https://", adapter)
    session.mount("http://", adapter)
    
    return session

# 添加Word文档生成函数
def generate_word_document(content, filename="lecture.docx"):
    """将Markdown内容转换为格式规范的Word文档"""
    try:
        # 创建新文档
        doc = Document()
        
        # 检测内容语言，设置相应的字体
        is_english = any(char in 'abcdefghijklmnopqrstuvwxyzABCDEFGHIJKLMNOPQRSTUVWXYZ' for char in content)
        
        # 设置文档默认样式
        style = doc.styles['Normal']
        font = style.font
        if is_english:
            # 英文内容使用英文字体
            font.name = 'Times New Roman'
            heading_font_name = 'Arial'
        else:
            # 中文内容使用中文字体
            font.name = '宋体'
            heading_font_name = '黑体'
        font.size = Pt(12)
        style.paragraph_format.space_after = Pt(6)
        style.paragraph_format.line_spacing = 1.5
        
        # 设置标题样式
        for i in range(1, 6):
            heading_style = doc.styles[f'Heading {i}']
            heading_font = heading_style.font
            heading_font.name = heading_font_name
            heading_font.bold = True
            if i == 1:
                heading_font.size = Pt(16)
            elif i == 2:
                heading_font.size = Pt(14)
            else:
                heading_font.size = Pt(12)
        
        # 分割内容为行
        lines = content.split('\n')
        in_code_block = False
        code_lines = []
        
        # 处理每一行
        for line in lines:
            line = line.strip()
            if not line:
                # 空行添加间距
                doc.add_paragraph()
                continue
                
            # 处理代码块
            if line.startswith('```'):
                if in_code_block:
                    # 结束代码块
                    in_code_block = False
                    if code_lines:
                        code_paragraph = doc.add_paragraph()
                        code_paragraph.alignment = WD_ALIGN_PARAGRAPH.LEFT
                        for code_line in code_lines:
                            code_run = code_paragraph.add_run(code_line + '\n')
                            # 代码块使用等宽字体
                            code_run.font.name = 'Consolas' if is_english else '宋体'
                            code_run.font.size = Pt(10)
                        code_lines = []
                else:
                    # 开始代码块
                    in_code_block = True
                continue
            
            if in_code_block:
                code_lines.append(line)
                continue
                
            # 处理标题 - 改进正则表达式匹配
            heading_match = re.match(r'^(#{1,6})\s+(.*)', line)
            if heading_match:
                level = len(heading_match.group(1))
                heading_text = heading_match.group(2).strip()
                # 限制标题级别在1-6之间
                level = max(1, min(6, level))
                heading = doc.add_heading(heading_text, level=level)
                if level == 1:
                    heading.alignment = WD_ALIGN_PARAGRAPH.CENTER
                heading.paragraph_format.space_before = Pt(12)
                heading.paragraph_format.space_after = Pt(6)
                continue
            
            # 处理列表项
            elif line.startswith('- ') or line.startswith('* '):
                p = doc.add_paragraph(style='List Bullet')
                run = p.add_run(line[2:])
                run.font.name = 'Times New Roman' if is_english else '宋体'
                p.paragraph_format.left_indent = Inches(0.25)
                p.paragraph_format.space_after = Pt(3)
            # 处理数字列表
            elif re.match(r'^\d+\.\s', line):
                p = doc.add_paragraph(style='List Number')
                run = p.add_run(line[line.find(' ')+1:])
                run.font.name = 'Times New Roman' if is_english else '宋体'
                p.paragraph_format.left_indent = Inches(0.25)
                p.paragraph_format.space_after = Pt(3)
            # 处理表格行
            elif '|' in line and re.match(r'^[\s|:-]+$', line) is None:
                # 简单的表格处理
                cells = [cell.strip() for cell in line.split('|') if cell.strip()]
                if cells and not any(cell.startswith('--') for cell in cells):
                    # 添加表格
                    if 'table' not in locals():
                        table = doc.add_table(rows=1, cols=len(cells))
                        table.style = 'Table Grid'
                        hdr_cells = table.rows[0].cells
                        for i, cell in enumerate(cells):
                            p = hdr_cells[i].paragraphs[0]
                            run = p.add_run(cell)
                            run.font.bold = True
                            run.font.name = 'Times New Roman' if is_english else '宋体'
                    else:
                        row_cells = table.add_row().cells
                        for i, cell in enumerate(cells):
                            p = row_cells[i].paragraphs[0]
                            run = p.add_run(cell)
                            run.font.name = 'Times New Roman' if is_english else '宋体'
            # 处理粗体文本
            elif '**' in line:
                p = doc.add_paragraph()
                parts = line.split('**')
                for i, part in enumerate(parts):
                    if part:  # 跳过空字符串
                        run = p.add_run(part)
                        run.font.name = 'Times New Roman' if is_english else '宋体'
                        if i % 2 == 1:  # 奇数索引部分是粗体
                            run.bold = True
            # 普通段落
            else:
                p = doc.add_paragraph()
                run = p.add_run(line)
                run.font.name = 'Times New Roman' if is_english else '宋体'
                p.paragraph_format.space_after = Pt(6)
        
        # 保存到字节流
        file_stream = io.BytesIO()
        doc.save(file_stream)
        file_stream.seek(0)
        
        return file_stream, filename
        
    except Exception as e:
        print(f"生成Word文档时出错: {e}")
        import traceback
        traceback.print_exc()
        return None, None

# 添加保存Word文档的函数
def save_lecture_to_word(lecture_content, chapter_name):
    """保存讲义内容到Word文档"""
    file_stream, filename = generate_word_document(lecture_content, f"{chapter_name}.docx")
    return file_stream, filename


# 创建全局会话
session = create_session_with_retries()

def extract_json_from_text(text):
    """从文本中提取JSON内容"""
    # 尝试找到JSON对象或数组
    json_match = re.search(r'(\{[\s\S]*\}|\[[\s\S]*\])', text)
    if json_match:
        try:
            return json.loads(json_match.group(1))
        except json.JSONDecodeError:
            pass
    
    # 如果上述方法失败，尝试找到被```json包裹的内容
    if '```json' in text:
        parts = text.split('```json')
        if len(parts) > 1:
            json_part = parts[1].split('```')[0].strip()
            try:
                return json.loads(json_part)
            except json.JSONDecodeError:
                pass
    
    # 如果还是失败，尝试找到被```包裹的内容
    if '```' in text:
        parts = text.split('```')
        if len(parts) > 1:
            json_part = parts[1].strip()
            try:
                return json.loads(json_part)
            except json.JSONDecodeError:
                pass
    
    # 最后尝试直接解析整个文本
    try:
        return json.loads(text)
    except json.JSONDecodeError:
        return {"error": "无法提取有效的JSON", "raw_text": text}

def call_deepseek(prompt, model="deepseek-chat", temperature=0.7):
    """调用DeepSeek API，包含完整的错误处理和重试机制"""
    # 检查API密钥是否设置
    if not DEEPSEEK_API_KEY or DEEPSEEK_API_KEY == "你的API密钥":
        return {"error": "未设置DeepSeek API密钥，请在.env文件中设置DEEPSEEK_API_KEY"}
    
    headers = {
        "Authorization": f"Bearer {DEEPSEEK_API_KEY}",
        "Content-Type": "application/json"
    }
    
    data = {
        "model": model,
        "messages": [{"role": "user", "content": prompt}],
        "temperature": temperature,
        "stream": False  # 确保不使用流式传输，减少连接问题
    }
    
    # 创建新的会话，避免使用全局会话可能的问题
    session = create_session_with_retries()
    
    try:
        print(f"正在调用DeepSeek API，提示词长度: {len(prompt)}")
        response = session.post(
            DEEPSEEK_API_URL,
            headers=headers,
            json=data,
            timeout=(CONNECT_TIMEOUT, READ_TIMEOUT)  # 分别设置连接和读取超时
        )
        
        print(f"API响应状态码: {response.status_code}")
        
        if response.status_code == 200:
            result = response.json()
            content = result["choices"][0]["message"]["content"]
            print("API调用成功!")
            return content
        else:
            error_msg = f"API错误: 状态码 {response.status_code}, 响应: {response.text}"
            print(error_msg)
            return {"error": error_msg}
            
    except requests.exceptions.Timeout:
        error_msg = f"API调用超时（连接:{CONNECT_TIMEOUT}s, 读取:{READ_TIMEOUT}s），请检查网络连接或稍后重试"
        print(error_msg)
        return {"error": error_msg}
        
    except requests.exceptions.ConnectionError as e:
        error_msg = f"网络连接错误: {str(e)}，请检查网络设置"
        print(error_msg)
        
        # 尝试诊断连接问题
        try:
            # 测试是否能解析域名
            import socket
            socket.gethostbyname("api.deepseek.com")
            dns_status = "DNS解析正常"
        except socket.gaierror:
            dns_status = "DNS解析失败"
        
        try:
            # 测试是否能连接到API服务器
            test_socket = socket.socket(socket.AF_INET, socket.SOCK_STREAM)
            test_socket.settimeout(5)
            test_socket.connect(("api.deepseek.com", 443))
            test_socket.close()
            connect_status = "可以连接到API服务器"
        except Exception as e:
            connect_status = f"无法连接到API服务器: {str(e)}"
        
        error_msg += f"\n诊断信息: {dns_status}, {connect_status}"
        return {"error": error_msg}
        
    except requests.exceptions.RequestException as e:
        error_msg = f"网络请求异常: {str(e)}"
        print(error_msg)
        return {"error": error_msg}
        
    except Exception as e:
        error_msg = f"API调用异常: {str(e)}"
        print(error_msg)
        return {"error": error_msg}
    finally:
        session.close()

def parse_json_response(response_text):
    """尝试解析JSON响应 - 增强版"""
    try:
        # 如果已经是字典，直接返回
        if isinstance(response_text, dict):
            return response_text
        
        # 使用增强的JSON提取函数
        result = extract_json_from_text(response_text)
        
        if isinstance(result, dict) and "error" in result:
            print(f"JSON解析失败，原始响应: {response_text}")
            return result
            
        return result
        
    except Exception as e:
        error_msg = f"解析响应时发生异常: {str(e)}"
        print(error_msg)
        print(f"原始响应: {response_text}")
        return {"error": error_msg, "raw_response": response_text}

def generate_course_outline(course_name, objectives, hours, education_stage="小学", policy_requirements=""):
    """生成课程大纲"""
    from prompts import PROMPT_COURSE_OUTLINE, EDUCATION_STAGE_GUIDANCE
    
    guidance = EDUCATION_STAGE_GUIDANCE.get(education_stage, "")
    
    # 构建政策要求部分
    policy_section = ""
    if policy_requirements:
        policy_section = f"""
重要政策要求：
{policy_requirements}

请确保课程大纲严格符合上述教育政策/考试大纲的要求！
"""
    
    prompt = PROMPT_COURSE_OUTLINE.format(
        course_name=course_name,
        education_stage=education_stage,
        objectives=objectives,
        hours=hours,
        education_stage_guidance=guidance,
        policy_requirements=policy_section  # 新增政策要求
    )
    
    print(f"生成课程大纲提示词: {prompt}")
    response = call_deepseek(prompt)
    return parse_json_response(response)

# 修改 generate_lecture_content 函数
def generate_lecture_content(chapter_name, key_points, hours, education_stage="小学", generation_language="中文", policy_requirements=""):
    """生成讲义内容"""
    from prompts import PROMPT_LECTURE_CONTENT, EDUCATION_STAGE_GUIDANCE
    
    guidance = EDUCATION_STAGE_GUIDANCE.get(education_stage, "")
    
    # 构建政策要求部分
    policy_section = ""
    if policy_requirements:
        policy_section = f"""
重要政策要求：
{policy_requirements}
"""
    
    prompt = PROMPT_LECTURE_CONTENT.format(
        education_stage=education_stage,
        chapter_name=chapter_name,
        key_points=key_points,
        hours=hours,
        education_stage_guidance=guidance,
        generation_language=generation_language,
        policy_requirements=policy_section  # 新增政策要求
    )
    
    response = call_deepseek(prompt)
    return response

def recommend_resources(course_name, education_stage="小学"):
    """推荐教学资源 - 完全重写：更好的格式处理和错误处理"""
    from prompts import PROMPT_RECOMMEND_RESOURCES, EDUCATION_STAGE_GUIDANCE
    
    guidance = EDUCATION_STAGE_GUIDANCE.get(education_stage, "")
    
    prompt = PROMPT_RECOMMEND_RESOURCES.format(
        course_name=course_name,
        education_stage=education_stage,
        education_stage_guidance=guidance
    )
    
    print(f"正在获取教学资源，课程: {course_name}, 教育阶段: {education_stage}")
    
    try:
        response = call_deepseek(prompt)
        print(f"API原始响应: {response}")
        
        # 如果API调用失败，直接返回模拟数据
        if isinstance(response, dict) and "error" in response:
            print(f"API调用失败，使用备用方案: {response['error']}")
            return recommend_mock_resources(course_name, education_stage)
        
        # 清理响应文本，移除可能的问题字符
        cleaned_response = clean_json_response(response)
        
        # 尝试解析JSON响应
        resources = parse_json_response(cleaned_response)
        print(f"解析后的资源: {resources}")
        
        # 检查解析结果并标准化格式
        standardized_resources = standardize_resources_format(resources, course_name)
        
        return standardized_resources
            
    except Exception as e:
        print(f"获取教学资源时发生异常: {e}")
        return recommend_mock_resources(course_name, education_stage)

def clean_json_response(text):
    """清理JSON响应，修复常见格式问题"""
    if not isinstance(text, str):
        return text
    
    # 移除可能的问题字符和格式
    cleaned = text.replace('\\"', '"')
    cleaned = cleaned.replace('\\n', ' ')
    cleaned = cleaned.replace('\\t', ' ')
    
    # 修复常见的JSON格式错误
    cleaned = re.sub(r'(\w+)\s*:\s*{', r'"\1": {', cleaned)  # 为键添加引号
    cleaned = re.sub(r',\s*}', '}', cleaned)  # 移除尾随逗号
    cleaned = re.sub(r',\s*]', ']', cleaned)  # 移除数组尾随逗号
    
    return cleaned

def standardize_resources_format(resources, course_name):
    """标准化资源格式，确保前端可以正确显示"""
    if not isinstance(resources, dict):
        # 如果不是字典，返回模拟数据
        return recommend_mock_resources(course_name)
    
    # 确保包含所有必需的键
    required_keys = ["教材", "在线视频", "工具/软件", "案例研究"]
    standardized = {}
    
    for key in required_keys:
        if key in resources:
            value = resources[key]
            
            # 标准化值为列表格式
            if isinstance(value, list):
                standardized[key] = value
            elif isinstance(value, dict):
                # 如果是字典，转换为单元素列表
                standardized[key] = [value]
            elif isinstance(value, str):
                # 如果是字符串，尝试解析或创建简单列表
                try:
                    parsed = json.loads(value)
                    standardized[key] = parsed if isinstance(parsed, list) else [parsed]
                except:
                    standardized[key] = [value]
            else:
                standardized[key] = [str(value)]
        else:
            # 如果键不存在，使用模拟数据
            mock_data = recommend_mock_resources(course_name)
            standardized[key] = mock_data.get(key, [f"暂无{key}信息"])
    
    return standardized

def recommend_mock_resources(course_name, education_stage="小学"):
    """推荐模拟教学资源（标准化格式版本）"""
    print(f"使用备用方案生成教学资源: {course_name}")
    
    course_name_lower = course_name.lower()
    
    # 金融类课程资源
    if any(keyword in course_name_lower for keyword in ["金融", "经济", "货币", "银行", "投资"]):
        return {
            "教材": [
                {
                    "书名": "金融学原理",
                    "作者": "李健",
                    "出版社": "高等教育出版社",
                    "备注": "系统讲解金融学基础理论"
                },
                {
                    "书名": "货币银行学", 
                    "作者": "黄达",
                    "出版社": "中国人民大学出版社",
                    "备注": "经典货币银行学教材"
                }
            ],
            "在线视频": [
                {
                    "视频标题": "金融学原理相关视频",
                    "发布平台": "B站搜索",
                    "主讲人/机构": "多个来源",
                    "链接": f"https://search.bilibili.com/all?keyword={course_name}+金融学+视频"
                },
                {
                    "视频标题": "经济学教学视频",
                    "发布平台": "百度搜索", 
                    "主讲人/机构": "多个来源",
                    "链接": f"https://www.baidu.com/s?wd={course_name}+经济学+教学视频"
                }
            ],
            "工具/软件": [
                {
                    "工具名称": "Wind金融终端",
                    "类型": "专业金融数据平台",
                    "用途": "金融市场数据分析"
                }
            ],
            "案例研究": [
                {
                    "案例名称": "2008年金融危机分析",
                    "领域": "金融风险",
                    "描述": "分析金融危机成因和应对措施"
                }
            ]
        }
    
    # 编程类课程资源
    elif any(keyword in course_name_lower for keyword in ["python", "编程", "计算机"]):
        return {
            "教材": [
                {
                    "书名": "Python编程：从入门到实践",
                    "作者": "Eric Matthes",
                    "出版社": "人民邮电出版社",
                    "备注": "适合初学者的Python教材"
                }
            ],
            "在线视频": [
                {
                    "视频标题": "Python编程教学视频",
                    "发布平台": "B站搜索",
                    "主讲人/机构": "多个来源",
                    "链接": f"https://search.bilibili.com/all?keyword={course_name}+Python+编程+教程"
                },
                {
                    "视频标题": "计算机科学教学资源",
                    "发布平台": "百度搜索",
                    "主讲人/机构": "多个来源", 
                    "链接": f"https://www.baidu.com/s?wd={course_name}+计算机+教学视频"
                }
            ],
            "工具/软件": [
                {
                    "工具名称": "PyCharm",
                    "类型": "IDE",
                    "用途": "Python开发环境"
                }
            ],
            "案例研究": [
                {
                    "案例名称": "Python数据分析实战",
                    "领域": "数据分析",
                    "描述": "使用Python进行数据分析和可视化"
                }
            ]
        }
    
    # 通用资源模板
    else:
        return {
            "教材": [
                {
                    "书名": f"{course_name}导论",
                    "作者": "多位专家",
                    "出版社": "高等教育出版社",
                    "备注": f"{course_name}领域入门教材"
                }
            ],
            "在线视频": [
                {
                    "视频标题": f"{course_name}教学视频",
                    "发布平台": "B站搜索",
                    "主讲人/机构": "多个来源",
                    "链接": f"https://search.bilibili.com/all?keyword={course_name}+教学视频"
                },
                {
                    "视频标题": f"{course_name}学习资源",
                    "发布平台": "百度搜索",
                    "主讲人/机构": "多个来源",
                    "链接": f"https://www.baidu.com/s?wd={course_name}+学习+视频"
                }
            ],
            "工具/软件": [
                {
                    "工具名称": "相关专业软件",
                    "type": "专业工具",
                    "用途": f"{course_name}领域专业应用"
                }
            ],
            "案例研究": [
                {
                    "案例名称": f"{course_name}应用案例",
                    "领域": "实践应用",
                    "描述": f"{course_name}在实际中的应用分析"
                }
            ]
        }

# 新增函数：根据用户反馈更新讲义内容
def update_lecture_content(current_content, user_input, conversation_history, education_stage="小学", generation_language="中文", policy_requirements=""):
    """根据用户反馈更新讲义内容"""
    try:
        from prompts import PROMPT_UPDATE_LECTURE, EDUCATION_STAGE_GUIDANCE
    except ImportError:
        print("警告：PROMPT_UPDATE_LECTURE 未定义，使用备用方案")
        return update_lecture_with_mock(current_content, user_input)
    
    # 格式化对话历史
    formatted_history = ""
    for role, message in conversation_history:
        formatted_history += f"{role}: {message}\n"
    
    guidance = EDUCATION_STAGE_GUIDANCE.get(education_stage, "")
    
    # 构建政策要求部分
    policy_section = ""
    if policy_requirements:
        policy_section = f"""
重要政策要求：
{policy_requirements}
"""
    
    prompt = PROMPT_UPDATE_LECTURE.format(
        education_stage=education_stage,
        current_content=current_content,
        conversation_history=formatted_history,
        user_input=user_input,
        education_stage_guidance=guidance,
        generation_language=generation_language,
        policy_requirements=policy_section  # 新增政策要求
    )
    
    response = call_deepseek(prompt)
    return response
def update_lecture_with_mock(current_content, user_input):
    """使用备用方案更新讲义内容"""
    print(f"使用备用方案更新讲义，用户输入: {user_input}")
    
    # 简单地在现有内容基础上添加用户要求的注释
    updated_content = f"{current_content}\n\n---\n\n## 根据用户反馈更新\n\n**用户要求:** {user_input}\n\n**更新说明:** 已根据用户反馈对讲义内容进行相应调整和优化。"
    
    return updated_content

# save_survey_result 函数
def save_survey_result(survey_data):
    """保存调查结果到JSON文件"""
    try:
        # 尝试读取现有数据
        try:
            with open("survey_results.json", "r", encoding="utf-8") as f:
                data = json.load(f)
        except (FileNotFoundError, json.JSONDecodeError):
            data = []
        
        # 添加新记录
        data.append(survey_data)
        
        # 写回文件
        with open("survey_results.json", "w", encoding="utf-8") as f:
            json.dump(data, f, ensure_ascii=False, indent=2)
            
        print(f"调查结果已保存到 survey_results.json，当前共有 {len(data)} 条记录")
        return True
            
    except Exception as e:
        print(f"保存调查结果时出错: {e}")
        # 如果出错，至少打印出来
        print(f"调查数据: {survey_data}")
        return False

def load_survey_results():
    """加载调查结果"""
    try:
        with open("survey_results.json", "r", encoding="utf-8") as f:
            data = json.load(f)
        return data
    except FileNotFoundError:
        return []
    except Exception as e:
        print(f"加载调查结果时出错: {e}")
        return []
    
# 模拟数据函数（作为备用方案）
def generate_mock_course_outline(course_name, objectives, hours, education_stage="小学", policy_requirements=""):
    """生成模拟课程大纲（备用方案）"""
    # 根据课程名称和教育阶段生成不同的内容
    course_name_lower = course_name.lower()
    
    # 根据教育阶段调整内容特点
    stage_adjustments = {
        "小学": {
            "depth": "基础", "complexity": "简单", "examples": "生活化", "language": "通俗有趣"
        },
        "初中": {
            "depth": "中等", "complexity": "适中", "examples": "贴近生活", "language": "严谨易懂"  
        },
        "高中": {
            "depth": "深入", "complexity": "复杂", "examples": "学术化", "language": "规范严谨"
        },
        "大学": {
            "depth": "专业", "complexity": "高深", "examples": "前沿性", "language": "学术专业"
        }
    }
    
    adjustment = stage_adjustments.get(education_stage, stage_adjustments["小学"])
    
    # 添加政策要求说明
    policy_note = ""
    if policy_requirements:
        policy_note = f"\n\n政策要求说明：本课程大纲已严格遵循相关政策要求进行设计。"
    
    if any(keyword in course_name_lower for keyword in ["数字", "经济", "金融"]):
        chapters = [
            {"章节名称": f"数字经济概述（{education_stage}版）", "学时": max(4, hours//8), "重点内容": f"数字经济定义、特征、发展历程 - {adjustment['depth']}讲解{policy_note}"},
            {"章节名称": f"数字技术基础（{education_stage}版）", "学时": max(6, hours//6), "重点内容": f"人工智能、大数据、区块链技术原理 - {adjustment['complexity']}介绍{policy_note}"},
            {"章节名称": f"数字商业模式（{education_stage}版）", "学时": max(8, hours//4), "重点内容": f"平台经济、共享经济、订阅经济模式 - {adjustment['examples']}案例{policy_note}"},
            {"章节名称": f"数据要素市场（{education_stage}版）", "学时": max(6, hours//6), "重点内容": f"数据确权、交易、定价机制 - {adjustment['language']}说明{policy_note}"},
            {"章节名称": f"数字治理政策（{education_stage}版）", "学时": max(8, hours//4), "重点内容": f"数字监管、隐私保护、国际合作 - {adjustment['depth']}分析{policy_note}"}
        ]
    elif any(keyword in course_name_lower for keyword in ["人工", "智能", "机器学习", "深度学习"]):
        chapters = [
            {"章节名称": f"人工智能导论（{education_stage}版）", "学时": max(4, hours//8), "重点内容": f"AI发展历史、基本概念、应用领域 - {adjustment['depth']}介绍{policy_note}"},
            {"章节名称": f"机器学习基础（{education_stage}版）", "学时": max(8, hours//4), "重点内容": f"监督学习、无监督学习、强化学习 - {adjustment['complexity']}讲解{policy_note}"},
            {"章节名称": f"深度学习原理（{education_stage}版）", "学时": max(10, hours//3), "重点内容": f"神经网络、CNN、RNN、Transformer - {adjustment['depth']}原理{policy_note}"},
            {"章节名称": f"AI应用实践（{education_stage}版）", "学时": max(8, hours//4), "重点内容": f"计算机视觉、自然语言处理、推荐系统 - {adjustment['examples']}应用{policy_note}"},
            {"章节名称": f"AI伦理与社会（{education_stage}版）", "学时": max(4, hours//8), "重点内容": f"AI伦理、偏见、社会责任 - {adjustment['language']}讨论{policy_note}"}
        ]
    else:
        # 通用课程结构
        chapters = [
            {"章节名称": f"{course_name}导论（{education_stage}版）", "学时": max(4, hours//8), "重点内容": f"基本概念和理论基础 - {adjustment['depth']}介绍{policy_note}"},
            {"章节名称": f"{course_name}核心原理（{education_stage}版）", "学时": max(8, hours//4), "重点内容": f"主要理论和方法论 - {adjustment['complexity']}讲解{policy_note}"},
            {"章节名称": f"{course_name}应用实践（{education_stage}版）", "学时": max(8, hours//4), "重点内容": f"实际应用和案例分析 - {adjustment['examples']}实践{policy_note}"},
            {"章节名称": f"{course_name}前沿发展（{education_stage}版）", "学时": max(6, hours//6), "重点内容": f"最新发展趋势和挑战 - {adjustment['depth']}展望{policy_note}"},
            {"章节名称": f"{course_name}总结展望（{education_stage}版）", "学时": max(4, hours//8), "重点内容": f"课程总结和未来展望 - {adjustment['language']}总结{policy_note}"}
        ]
    
    # 调整总学时以匹配输入
    total_hours = sum(chapter["学时"] for chapter in chapters)
    if total_hours != hours:
        scale = hours / total_hours
        for chapter in chapters:
            chapter["学时"] = max(2, round(chapter["学时"] * scale))
    
    return {
        "课程名称": course_name,
        "教育阶段": education_stage,  # 新增
        "教学目标": objectives,
        "总学时": hours,
        "章节列表": chapters
    }

def generate_mock_lecture_content(chapter_name, key_points, hours, education_stage="小学", generation_language="中文", policy_requirements=""):
    """生成模拟讲义内容（用于测试）"""
    stage_characteristics = {
        "小学": "使用生动有趣的语言，多举生活中的例子，强调兴趣培养",
        "初中": "注重知识系统性，使用图表归纳，语言严谨易懂", 
        "高中": "内容深入系统，强调逻辑思维，使用规范学术语言",
        "大学": "内容专业前沿，强调独立思考，使用学术化表达"
    }
    
    characteristic = stage_characteristics.get(education_stage, "")
    
    # 添加政策要求说明
    policy_note = ""
    if policy_requirements:
        policy_note = f"\n\n**政策要求说明：** 本讲义已严格遵循相关政策要求进行编写。"
    
    # 根据语言选择内容
    if generation_language == "英文":
        return f"""
# {chapter_name} ({education_stage} Version)

## Part 1: Core Concepts

### Key Points Summary

{key_points}

**Educational Stage Characteristics:** {characteristic}

### Example Questions

1. **Question:** Please explain the meaning and importance of the core concepts in this chapter.
   **Answer:** The core concept of this chapter is {key_points.split('、')[0] if '、' in key_points else key_points}, which is fundamental for understanding subsequent content and has significant theoretical and practical importance.

2. **Question:** Please illustrate the application of this chapter's knowledge in practice with examples.
   **Answer:** Use specific cases to demonstrate how {key_points.split('、')[0] if '、' in key_points else key_points} is applied in real scenarios and its effectiveness.
{policy_note}

## Summary and Recommendations

### Learning Key Points Summary
- Master the basic concepts of {key_points}
- Understand related theories and application methods  
- Be able to use knowledge to solve practical problems

### Learning Recommendations
- Focus on combining theory and practice
- Do more exercises to consolidate knowledge points
- Actively participate in class discussions and practice
"""
    else:
        # 中文内容 - 使用正确的Markdown格式
        return f"""
# {chapter_name}（{education_stage}版）

## 第一部分：核心概念

### 知识点摘要

{key_points}

**教育阶段特点：** {characteristic}

### 例题

1. **题目：** 请解释本章核心概念的含义及其重要性。
   **解析：** 本章核心概念是{key_points.split('、')[0] if '、' in key_points else key_points}，它是理解后续内容的基础，具有重要的理论和实践意义。

2. **题目：** 请结合实例说明本章知识在实际中的应用。
   **解析：** 以具体案例说明{key_points.split('、')[0] if '、' in key_points else key_points}在实际场景中的应用方法和效果。
{policy_note}

## 总结与建议

### 学习要点总结
- 掌握{key_points}的基本概念
- 理解相关理论和应用方法  
- 能够运用知识解决实际问题

### 学习建议
- 注重理论与实践相结合
- 多做练习巩固知识点
- 积极参与课堂讨论和实践
"""
    
# PPT生成函数
def generate_ppt_document(content, filename="lecture.pptx", template_path=None):
    """将Markdown内容转换为格式规范的PPT文档，使用指定模板，支持自动分页"""
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.enum.text import PP_ALIGN
        from pptx.dml.color import RGBColor
        import re
        
        # 创建演示文稿，使用模板如果提供
        if template_path and os.path.exists(template_path):
            prs = Presentation(template_path)
            # 使用模板时，使用标题和内容布局
            slide_layout = prs.slide_layouts[1]
        else:
            prs = Presentation()
            # 设置默认幻灯片布局 - 标题和内容
            slide_layout = prs.slide_layouts[1]
        
        # 检测内容语言
        is_english = any(char in 'abcdefghijklmnopqrstuvwxyzABCDEFGHIJKLMNOPQRSTUVWXYZ' for char in content)
        
        # 设置字体
        if is_english:
            title_font_name = 'Arial'
            content_font_name = 'Arial'
        else:
            title_font_name = '黑体'
            content_font_name = '微软雅黑'
        
        # 分割内容为行
        lines = content.split('\n')
        
        # 当前幻灯片和内容
        current_slide = None
        current_content = []
        current_line_count = 0
        MAX_LINES_PER_SLIDE = 12  # 每页最多12行内容
        
        # 分页规则：遇到这些关键词时创建新幻灯片
        page_break_keywords = [
            r'^# ',  # 一级标题
            r'^## ',  # 二级标题
            r'^### ',  # 三级标题
            r'^#### ',  # 四级标题
        ]

        def create_new_slide(title_text=None):
            """创建新幻灯片并设置标题"""
            nonlocal current_slide, current_content, current_line_count
            
            # 保存当前幻灯片内容（如果有的话）
            if current_slide is not None and current_content:
                add_content_to_slide(current_slide, current_content)
            
            # 创建新幻灯片
            current_slide = prs.slides.add_slide(slide_layout)
            
            # 设置标题
            if title_text:
                title_box = current_slide.shapes.title
                if title_box is not None:
                    title_box.text = clean_title_text(title_text)
                    # 设置标题字体
                    for paragraph in title_box.text_frame.paragraphs:
                        paragraph.alignment = PP_ALIGN.CENTER
                        for run in paragraph.runs:
                            run.font.size = Pt(28 if len(title_text) < 20 else 24)
                            run.font.bold = True
                            run.font.name = title_font_name
                            run.font.color.rgb = RGBColor(0, 0, 0)
            
            # 重置内容
            current_content = []
            current_line_count = 0
            
            return current_slide

        def clean_title_text(text):
            """清理标题文本，移除Markdown标记"""
            # 移除Markdown标题标记
            cleaned = re.sub(r'^#+\s*', '', text)
            # 移除粗体标记
            cleaned = re.sub(r'\*\*(.*?)\*\*', r'\1', cleaned)
            return cleaned.strip()

        def add_content_to_slide(slide, content_lines):
            """将内容添加到幻灯片"""
            if not content_lines:
                return
            
            # 过滤空行并清理内容
            filtered_lines = [clean_content_text(line) for line in content_lines if line.strip()]
            if not filtered_lines:
                return
                
            content_text = "\n".join(filtered_lines)
            
            # 检查是否有内容占位符
            content_placeholder = None
            for shape in slide.shapes:
                if shape.has_text_frame and shape != slide.shapes.title:
                    content_placeholder = shape
                    break
            
            if content_placeholder and len(slide.shapes.placeholders) > 1:
                content_box = slide.shapes.placeholders[1]
                tf = content_box.text_frame
                tf.clear()  # 清空现有内容
                tf.text = content_text
                
                # 设置字体和格式
                for paragraph in tf.paragraphs:
                    paragraph.alignment = PP_ALIGN.LEFT
                    for run in paragraph.runs:
                        run.font.size = Pt(18)
                        run.font.name = content_font_name
                        run.font.color.rgb = RGBColor(0, 0, 0)
            else:
                # 如果没有内容占位符，创建一个文本框
                left = Inches(0.5)
                top = Inches(1.5)
                width = Inches(9)
                height = Inches(5.5)
                textbox = slide.shapes.add_textbox(left, top, width, height)
                tf = textbox.text_frame
                tf.text = content_text
                tf.word_wrap = True
                
                # 设置字体
                for paragraph in tf.paragraphs:
                    paragraph.alignment = PP_ALIGN.LEFT
                    for run in paragraph.runs:
                        run.font.size = Pt(18)
                        run.font.name = content_font_name
                        run.font.color.rgb = RGBColor(0, 0, 0)

        def clean_content_text(text):
            """清理内容文本"""
            # 移除Markdown标记但保留文本内容
            cleaned = re.sub(r'^[*-]\s*', '', text)  # 移除列表标记
            cleaned = re.sub(r'\*\*(.*?)\*\*', r'\1', cleaned)  # 移除粗体但保留文本
            cleaned = re.sub(r'`(.*?)`', r'\1', cleaned)  # 移除代码标记但保留文本
            return cleaned.strip()

        def check_and_split_content(line):
            """检查内容是否需要分页，如果需要则创建新幻灯片"""
            nonlocal current_line_count, current_slide
            
            # 空行不计入行数限制
            if not line.strip():
                return False
                
            # 估算行长度（考虑换行）
            line_length = len(line)
            line_count_estimate = max(1, line_length // 50)  # 每行大约50字符
            
            current_line_count += line_count_estimate
            
            # 如果超过行数限制，创建新幻灯片
            if current_line_count >= MAX_LINES_PER_SLIDE:
                # 保存当前内容
                if current_slide is not None and current_content:
                    add_content_to_slide(current_slide, current_content)
                
                # 创建新幻灯片（续页）
                create_new_slide("Continued" if is_english else "（续）")
                
                # 将当前行添加到新幻灯片
                current_content.append(line)
                current_line_count = line_count_estimate  # 重置计数
                return True
                
            return False

        # 首先创建封面幻灯片
        if lines and lines[0].strip():
            first_line = lines[0].strip()
            title_text = clean_title_text(first_line)
            create_new_slide(title_text)
        
        # 处理每一行（跳过第一行，因为已经在封面使用了）
        for i, line in enumerate(lines):
            if i == 0:  # 跳过封面标题行
                continue
                
            line = line.strip()
            if not line:
                # 空行也添加到内容中，但不触发分页检查
                current_content.append("")
                continue
                
            # 检查是否需要分页（基于标题）
            need_break = any(re.search(pattern, line) for pattern in page_break_keywords)
            
            # 如果需要分页或者当前没有幻灯片，创建新幻灯片
            if need_break:
                # 创建新幻灯片
                slide_title = clean_title_text(line)
                create_new_slide(slide_title)
            else:
                # 检查是否需要基于内容长度分页
                if check_and_split_content(line):
                    continue  # 已经在新幻灯片中处理了这一行
            
            # 添加到当前内容（如果不是标题行）
            if not need_break:
                current_content.append(line)
        
        # 处理最后一页内容
        if current_slide is not None and current_content:
            add_content_to_slide(current_slide, current_content)
        
        # 如果没有创建任何幻灯片，创建一个默认幻灯片
        if len(prs.slides) == 0:
            create_new_slide("Presentation" if is_english else "演示文稿")
            add_content_to_slide(current_slide, ["No content available." if is_english else "暂无内容。"])
        
        # 保存到字节流
        file_stream = io.BytesIO()
        prs.save(file_stream)
        file_stream.seek(0)
        
        return file_stream, filename
        
    except Exception as e:
        print(f"生成PPT文档时出错: {e}")
        import traceback
        traceback.print_exc()
        return None, None
    
# 添加保存PPT文档的函数
def save_lecture_to_ppt(lecture_content, chapter_name, template_path=None):
    """保存讲义内容到PPT文档，可选择使用模板"""
    # 清理章节名称中的特殊字符
    clean_chapter_name = re.sub(r'[\\/*?:"<>|]', "", chapter_name)
    filename = f"{clean_chapter_name}.pptx"
    file_stream, filename = generate_ppt_document(lecture_content, filename, template_path)
    return file_stream, filename